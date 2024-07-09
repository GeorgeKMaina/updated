# -*- coding: utf-8 -*-
"""
Created on Thu Jul  4 14:30:49 2024

@author: Gmaina
"""
import numpy as np
import pandas as pd
import streamlit as st
import datetime as dt
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import ChartData
#import chat_api_key
import openai
import os

def variable_inputs():
    sample_of_surveys = ('Course Evaluations',
                         'Satisfaction Surveys',
                         'First-Year Enrolment Survey',
                         'Continuing Students Survey',
                         'Graduate Exit Survey',
                         'Market Interest Survey',
                         'Event Feedback Survey',
                         'Other')

    name_of_survey = st.selectbox('What survey is this?', sample_of_surveys)
    if name_of_survey == 'Other':
        name_of_survey = st.text_area(
            label='Name of the survey',
            height=75,
            max_chars=50,
            placeholder='Write here'
        )

    # Upload the dataset
    dataset = st.file_uploader(label='Upload the dataset')

    # Check if the dataset is uploaded
    if dataset:
        # Determine the file type and load the dataset accordingly
        if dataset.name.endswith('.xlsx'):
            dataset = pd.read_excel(dataset)
        elif dataset.name.endswith('.csv'):
            dataset = pd.read_csv(dataset)
        else:
            st.error('Dataset should either be in CSV or Excel format')
        
        # Display the first few rows of the dataset
        st.write(dataset.head())
        
        # Number of responses
        st.write('Number of respondents', dataset.shape[0])
        
    else:
        st.info('Please upload a dataset')

    target_population = st.text_input("What is the target population: ")
    start_date = st.date_input("When was the survey launched?", value=dt.date(2024, 7, 4))
    end_date = st.date_input("When was the survey closed?", value=dt.date.today())
            
    # Upload the logo
    logo = st.file_uploader(label='Upload the institutional logo', type=["png", "jpg", "jpeg"])
    
    if logo:
        st.image(logo, width=200)

    # Optionally return the collected inputs for further processing
    return name_of_survey, dataset, target_population, start_date, end_date, logo

def generate_powerpoint(df, title_text='', logo_file='', executive_summary='', output_path='survey_analysis.pptx'):
    """
    Generate a PowerPoint presentation with slides based on survey data.

    Parameters:
        df (DataFrame): Survey data in DataFrame format.
        title_text (str): Title for the PowerPoint presentation.
        logo_file (str): File path of the logo image to be included in the slides.
        executive_summary (str): The executive summary text to be included in the presentation.
        output_path (str): File path to save the PowerPoint presentation.
    """
    # Create a new presentation with widescreen (16:9) slide size
    #https://python-pptx.readthedocs.io/en/latest/
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Add an executive summary slide
    exec_summary_slide_layout = prs.slide_layouts[1]  # Title and Content layout
    exec_summary_slide = prs.slides.add_slide(exec_summary_slide_layout)
    title = exec_summary_slide.shapes.title
    title.text = "Executive Summary"

    content_shape = exec_summary_slide.placeholders[1]
    text_frame = content_shape.text_frame
    p = text_frame.add_paragraph()
    p.text = executive_summary
    p.font.size = Pt(14)

    # Add the company logo to the executive summary slide
    if logo_file:
        exec_summary_slide.shapes.add_picture(logo_file, Inches(11.71), Inches(0.03), height=Inches(0.75))

    # Add a title slide
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = title_slide.shapes.title
    title.text = title_text

    # Add the company logo to every slide
    if logo_file:
        for slide in prs.slides:
            slide.shapes.add_picture(logo_file, Inches(11.71), Inches(0.03), height=Inches(0.75))

    # Iterate over each column in the DataFrame
    for column in df.columns:
        unique_values = df[column].nunique(dropna=True)
        slide_layout = prs.slide_layouts[5] if unique_values <= 2 else prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)

        # Add title to the slide
        title_shape = slide.shapes.title
        if not title_shape:
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_shape.text = column

        # Calculate percentages of responses for each category and round to nearest whole number
        total_responses = df[column].count()
        percentages = (df[column].value_counts(normalize=True) * 100).round()

        # Add content to the slide (pie chart or bar chart)
        chart_data = ChartData()
        chart_data.categories = percentages.index.tolist()
        chart_data.add_series(column, percentages.tolist())

        if unique_values <= 2:
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, Inches(0.5), Inches(1.5), Inches(8), Inches(5), chart_data
            ).chart
        else:
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.5), Inches(8), Inches(5), chart_data
            ).chart

        # Set chart properties
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False  # Ensure legend doesn't overlap chart
        chart.series[0].data_labels.show_percentage = True  # Show data labels as percentages
        chart.series[0].data_labels.number_format = '0%'  # Format data labels as percentages with "%" sign
        chart.series[0].data_labels.font.size = Pt(10)  # Set font size of data labels
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = RGBColor(52, 79, 158)

        # Set the font size of the chart title
        chart_title = chart.chart_title.text_frame
        chart_title.paragraphs[0].font.size = Pt(14)

        # Generate narrative for the current column using OpenAI
        narrative_prompt = (
            f"Generate a brief narrative for a chart based on the following data:\n"
            f"Column: {column}\n"
            f"Data: {percentages.to_dict()}\n"
            f"Provide a short explanation summarizing the key insights from this chart."
        )

        narrative_response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {
                    "role": "user",
                    "content": narrative_prompt,
                }
            ],
            max_tokens=150,
            n=1,
            stop=None,
            temperature=0.7,
        )

        narrative_text = narrative_response.choices[0].message['content'].strip()

        # Add the narrative to the slide
        left = Inches(0.5)
        top = Inches(6.0)
        width = Inches(12.0)
        height = Inches(1.0)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = narrative_text
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(12)

    # Save the presentation
    prs.save(output_path)
    os.startfile(output_path)

# Main function to run the Streamlit app
def main():
    st.title("Survey Analysis Automation")
    name_of_survey, dataset, target_population, start_date, end_date, logo = variable_inputs()

    # You can now use the collected inputs for further processing
    if dataset is not None:
        st.write(f"Survey: {name_of_survey}")
        st.write(f"Target Population: {target_population}")
        st.write(f"Survey Period: {start_date} to {end_date}")

        # Add a button to generate the PowerPoint
        if st.button('Generate PowerPoint'):
            # Set your OpenAI API key
            import os
            from dotenv import load_dotenv
            import openai
            
            # Load environment variables from the .env file
            load_dotenv()

            # Get the OpenAI API key from the environment variables
            openai.api_key=os.getenv("new_api_key")

            # Assign the API key to the OpenAI library
            openai.api_key 
            

            # Generate objectives based on the survey name
            objectives_prompt = (
                f"Based on the survey name '{name_of_survey}', generate a list of 3-5 objectives "
                "that this survey might aim to achieve. The objectives should be relevant to "
                "the survey context and provide insights into what the survey is trying to accomplish."
            )

            objectives_response = openai.ChatCompletion.create(
                model="gpt-3.5-turbo",
                messages=[
                    {
                        "role": "user",
                        "content": objectives_prompt,
                    }
                ],
                max_tokens=150,
                n=1,
                stop=None,
                temperature=0.7,
            )

            objectives = objectives_response.choices[0].message['content'].strip()

            executive_summary = (
                f"The survey titled '{name_of_survey}' was conducted from {start_date} to {end_date} "
                f"with a target population of {target_population}. The survey received a total of "
                f"{dataset.shape[0]} responses. The main objectives of the survey were:\n{objectives}"
            )

            # Save the logo file if it exists
            logo_path = None
            if logo:
                logo_path = "uploaded_logo.png"
                with open(logo_path, "wb") as f:
                    f.write(logo.getbuffer())

            # Generate the PowerPoint presentation
            generate_powerpoint(
                df=dataset,
                title_text=f"Survey Analysis: {name_of_survey}",
                logo_file=logo_path,
                executive_summary=executive_summary,
                output_path="survey_analysis.pptx"
            )

if __name__ == '__main__':
    main()
